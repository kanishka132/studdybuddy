import os
import re
import tempfile
from flask import Flask, request, jsonify
from flask_cors import CORS
from dotenv import load_dotenv
from werkzeug.utils import secure_filename
from supabase import create_client
from PyPDF2 import PdfReader
import google.generativeai as genai
import fitz
import docx2txt
import pptx

# === Load Environment Variables ===
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")

# === Configure Gemini Model ===
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel(model_name="gemini-1.5-flash")

# === Initialize Supabase ===
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)

# === Flask App Setup ===
app = Flask(__name__)
CORS(app)
app.config['UPLOAD_FOLDER'] = tempfile.gettempdir()

# === Helper Functions ===
def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        page_text = page.get_text("text")
        text += page_text
        print(f"📄 Page extracted: {len(page_text)} chars")

    doc.close()
    return text

def extract_text_from_docx(file_path):
    return docx2txt.process(file_path)

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path):
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".docx"):
        return extract_text_from_docx(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    return ""

def extract_text_from_supabase_paths(file_paths):
    all_text = ""
    for path in file_paths:
        ext = os.path.splitext(path)[1].lower()
        file_bytes = supabase.storage.from_("materials").download(path)
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            extracted = extract_text(tmp.name)
            print(f"✅ Extracted from {path}:\n", extracted[:10])  # show first 1000 characters

            all_text += extracted + "\n"
    return re.sub(r'\s+', ' ', all_text).strip()

def generate_ai_response(task, text, count=5, difficulty="medium"):
    if task == "summary":
        prompt = f"Please summarize the following content clearly and concisely:\n{text}"
    elif task == "quiz":
        prompt = f"""
        Create {count} {difficulty}-level multiple-choice questions based on the text below.
        Each should have 4 options (A, B, C, D) with one correct answer. Return JSON array:
        [
          {{
            "question": "...",
            "options": ["A", "B", "C", "D"],
            "correct_answer": 1
          }}
        ]\nText:\n{text}"""
    elif task == "flashcards":
        prompt = f"""
        Create {count} flashcards from the following text in JSON format:
        [
          {{
            "front": "Term or Question",
            "back": "Answer or Explanation"
          }}
        ]\nText:\n{text}"""
    else:
        return None

    response = model.generate_content([prompt])
    print("🧠 Gemini raw response:", response.text)
    return response.text

# === Routes ===
@app.route("/generate-learning-content", methods=["POST"])
def generate_learning_content():
    try:
        data = request.get_json()
        file_paths = data.get("file_paths", [])
        actions = data.get("actions", [])
        project_name = data.get("project_name", "Untitled Project")
        question_count = int(data.get("question_count", 5))
        difficulty = data.get("difficulty", "medium")

        text = extract_text_from_supabase_paths(file_paths)

        if not text.strip():
            print("⚠️ Extracted text is empty. Skipping AI generation.")
            return jsonify({
                "success": True,
                "project_name": project_name,
                "results": {}
            })


        results = {}
        if "summary" in actions:
            results["summary"] = generate_ai_response("summary", text)
        if "quiz" in actions:
            results["quiz"] = generate_ai_response("quiz", text, count=question_count, difficulty=difficulty)
        if "flashcards" in actions:
            results["flashcards"] = generate_ai_response("flashcards", text)

        print("text:", text[:100])
        print("results: ",results)

        return jsonify({
            "success": True,
            "project_name": project_name,
            "results": results
        })

    except Exception as e:
        print("Error in /generate-learning-content:", str(e))
        return jsonify({"success": False, "error": "Processing failed"}), 500

@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "healthy"})

# === Run App ===
if __name__ == "__main__":
    port = int(os.getenv("PORT", 5000))
    app.run(debug=True, host="0.0.0.0", port=port)
